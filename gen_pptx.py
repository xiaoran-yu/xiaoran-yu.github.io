import os, io
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Fix Chinese font
font_path = r'C:\Windows\Fonts\simkai.ttf'
if not os.path.exists(font_path):
    font_path = r'C:\Windows\Fonts\msyh.ttc'
if os.path.exists(font_path):
    fm.fontManager.addfont(font_path)
    plt.rcParams['font.family'] = fm.FontProperties(fname=font_path).get_name()
plt.rcParams['axes.unicode_minus'] = False

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ===== COLOR PALETTE =====
DARK_BLUE = RGBColor(0x1E, 0x3F, 0x73)
PURPLE    = RGBColor(0x5B, 0x3A, 0x6E)
GOLD      = RGBColor(0xB8, 0x94, 0x2E)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
BLACK     = RGBColor(0x1A, 0x20, 0x2C)
GRAY      = RGBColor(0x64, 0x74, 0x8B)
LIGHT_BG  = RGBColor(0xF4, 0xF6, 0xF8)
POS_GREEN = RGBColor(0x1A, 0x6B, 0x4A)
NEG_RED   = RGBColor(0xB8, 0x3A, 0x3A)
PALE_PURPLE = RGBColor(0xD4, 0xC5, 0xE0)
LIGHT_GRAY_BG = RGBColor(0xF8, 0xFA, 0xFC)

# Pie chart colors matching HTML
PIE_COLORS = ['#1E3F73','#2B5C9E','#3A7CC3','#4E94D6','#6AA8E0','#87BCE8','#A3CFEF','#B8D9F4','#C5E0F7','#D2E8FA','#E0F0FC','#E8F4FD','#F0F8FE','#F5FBFE']

def add_header(slide, title_text, page_num):
    hdr = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.88))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = DARK_BLUE; hdr.line.fill.background()
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.88), prs.slide_width, Inches(0.04))
    line.fill.solid(); line.fill.fore_color.rgb = GOLD; line.line.fill.background()
    tf = hdr.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
    p.text = title_text; p.font.size = Pt(21); p.font.bold = True; p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT; tf.margin_left = Inches(0.6); tf.margin_top = Inches(0.1)
    ft = slide.shapes.add_textbox(Inches(0.6), Inches(7.08), Inches(12), Inches(0.25))
    fp = ft.text_frame.paragraphs[0]; fp.text = "内部资料 · 仅供董事会及管理层使用"; fp.font.size = Pt(9); fp.font.color.rgb = GRAY
    fp2 = ft.text_frame.add_paragraph(); fp2.text = f"2026年6月  ·  {page_num}/3"; fp2.font.size = Pt(9); fp2.font.color.rgb = GRAY; fp2.alignment = PP_ALIGN.RIGHT

def set_cell(tbl, r, c, text, font_size=10, bold=False, color=BLACK, bg=None, align=PP_ALIGN.CENTER, font_name='Microsoft YaHei'):
    cell = tbl.cell(r, c); cell.text = ""
    p = cell.text_frame.paragraphs[0]; p.text = str(text); p.alignment = align
    p.font.size = Pt(font_size); p.font.bold = bold; p.font.color.rgb = color; p.font.name = font_name
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg: cell.fill.solid(); cell.fill.fore_color.rgb = bg

def fmt(n):
    if n is None or n == 0: return '—'
    return f'{int(n):,}'

def pfmt(n):
    s = f'{int(n):,}'
    return f'+{s}' if n >= 0 else s

def generate_pie_chart(data, title, filename, is_profit=False):
    """Generate pie chart matching HTML styling"""
    fig, ax = plt.subplots(figsize=(3.8, 3.2), facecolor='white')
    # Sort by value descending
    sorted_data = sorted(data, key=lambda x: abs(x[1]), reverse=True)
    labels = [d[0] for d in sorted_data]
    values = [abs(d[1]) for d in sorted_data]
    total = sum(values)

    colors = PIE_COLORS[:len(values)]
    wedges, texts, autotexts = ax.pie(values, labels=None, colors=colors,
                                        autopct=lambda pct: f'{pct:.1f}%' if pct > 5 else '',
                                        startangle=90, pctdistance=0.58,
                                        wedgeprops=dict(width=1, edgecolor='white', linewidth=0.8))

    for at in autotexts:
        at.set_fontsize(8); at.set_fontweight('bold'); at.set_color('white')
        at.set_fontfamily('sans-serif')

    ax.set_title(title, fontsize=11, fontweight='bold', color='#1E3F73', pad=8)

    buf = io.BytesIO()
    plt.tight_layout()
    plt.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor='white', edgecolor='none')
    plt.close()
    buf.seek(0)
    return buf

# ===== DATA =====
bus_data = [
    ('PE', 11923, 17852, 602, 6122), ('RE', 30200, 7508, 9330, -12541),
    ('弘源(RE)', 34495, 31373, 7547, 9982), ('绿创', 9804, 14208, 3573, 5136),
    ('大运河', 6710, 8545, 100, 753), ('西安', 2300, 2225, 20, -21),
    ('重庆', 2853, 0, -1912, 0), ('大健康', 16024, 6551, 4632, -3217),
    ('消费', 7276, 3793, 565, -2924), ('新材料', 7082, 3454, 2374, -770),
    ('半导体', 14025, 0, 4727, -6311), ('科创赋能', 9850, 5854, 1681, -1190),
    ('HV', 12649, 12842, 9165, 9343), ('弘毅创投', 2862, 0, -3534, -1046),
]

segment_data = [
    ('一级市场', '一级市场各BU', 168053, 114204, 36870, 3317, True),
    ('二级市场及恒盛', '远方公募', 5704, 3849, -43542, -36048, False),
    ('二级市场及恒盛', '金涌 (60.81%)', 64483, 50403, 41226, 30893, False),
    ('二级市场及恒盛', '恒盛 (56.25%)', 49108, 51211, -4496, -22446, True),
    ('集团本部', '本部各职能部门', 67612, 97701, -55027, -30083, False),
]

expense_data = [
    ('资讯系统', 4166, 1117, True), ('差旅', 2964, 2644, False), ('会议费', 360, 743, False),
    ('交通费', 350, 243, False), ('培训招聘', 208, 542, False), ('人力资源费', 15353, 12266, False),
    ('市场发展费', 4668, 4997, False), ('业务招待费', 1520, 1407, False), ('员工福利费', 787, 1036, False),
    ('其他业务支出', 13028, 12479, False), ('专业咨询费', 8252, 12968, False),
]

# ============ SLIDE 1: 集团整体情况 ============
s1 = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s1, "2026年度集团预算分析 — 集团整体情况", "1")

t1 = s1.shapes.add_table(12, 7, Inches(0.6), Inches(1.15), Inches(12.1), Inches(5.6)).table

h1 = ['板块', '收入\n2026预算', '收入\n2025实际', '收入\n变动额', '分摊后利润\n2026预算', '分摊后利润\n2025实际', '利润\n变动额']
for c, h in enumerate(h1):
    bg = PURPLE if c >= 4 else DARK_BLUE
    set_cell(t1, 0, c, h, 11, True, WHITE, bg)

ri = 1; sR26 = sR25 = sP26 = sP25 = 0; last_g = ''
for gname, name, r26, r25, p26, p25, gold in segment_data:
    if gname != last_g:
        last_g = gname
        for c in range(7): set_cell(t1, ri, c, gname, 12, True, DARK_BLUE, LIGHT_BG, PP_ALIGN.LEFT)
        ri += 1
    sR26 += r26; sR25 += r25; sP26 += p26; sP25 += p25
    rc, pc = r26 - r25, p26 - p25
    row_bg = LIGHT_BG if gold else None
    set_cell(t1, ri, 0, name, 13, True, BLACK, row_bg, PP_ALIGN.LEFT)
    vals = [(fmt(r26), BLACK, 13), (fmt(r25), BLACK, 13), (pfmt(rc), POS_GREEN if rc>=0 else NEG_RED, 13),
            (fmt(p26), POS_GREEN if p26>=0 else NEG_RED, 13), (fmt(p25), POS_GREEN if p25>=0 else NEG_RED, 13),
            (pfmt(pc), POS_GREEN if pc>=0 else NEG_RED, 13)]
    for ci, (v, clr, fs) in enumerate(vals):
        set_cell(t1, ri, ci+1, v, fs, False, clr, row_bg)
    ri += 1

gr, gp = sR26 - sR25, sP26 - sP25
for c, v in enumerate(['集团总合计', sR26, sR25, gr, sP26, sP25, gp]):
    clr = BLACK if c == 0 else (BLACK if c < 3 else BLACK)
    set_cell(t1, ri, c, fmt(v) if isinstance(v, int) else v, 14, True, BLACK, LIGHT_BG, PP_ALIGN.LEFT if c==0 else PP_ALIGN.CENTER)

fn1 = s1.shapes.add_textbox(Inches(0.6), Inches(6.85), Inches(12), Inches(0.3))
fn1.text_frame.word_wrap = True
p = fn1.text_frame.paragraphs[0]; p.font.size = Pt(10); p.font.color.rgb = GRAY
p.text = "注：金涌持股比例60.81%，恒盛持股比例56.25%。总部各职能部门收入包括PE部门按85%核算的管理费收入人民币55,975K，以及投资管理部2026年预计收入人民币9,707K。"

# ============ SLIDE 2: 一级市场各BU + 饼图 ============
s2 = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s2, "2026年度集团预算分析 — 一级市场各BU", "2")

# LEFT: Table (~62% width)
t2 = s2.shapes.add_table(17, 9, Inches(0.35), Inches(1.1), Inches(7.8), Inches(5.7)).table

h2 = ['BU', '收入\n2026预算', '收入\n2025实际', '变动额', '变动率', '分摊后利润\n2026预算', '分摊后利润\n2025实际', '变动额', '变动率']
for c, h in enumerate(h2):
    bg = PURPLE if c >= 5 else DARK_BLUE
    set_cell(t2, 0, c, h, 9, True, WHITE, bg)

R26 = R25 = P26 = P25 = 0
for i, (name, rev26, rev25, prof26, prof25) in enumerate(bus_data):
    R26 += rev26; R25 += rev25; P26 += prof26; P25 += prof25
    rc, pc = rev26 - rev25, prof26 - prof25
    rc_pct = f"{'+' if rc>=0 else ''}{rc/rev25*100:.0f}%" if rev25 else '—'
    pc_pct = f"{'+' if pc>=0 else ''}{pc/abs(prof25)*100:.0f}%" if prof25 else '—'
    ri = i + 1
    set_cell(t2, ri, 0, name, 10, True, BLACK, None, PP_ALIGN.LEFT)
    set_cell(t2, ri, 1, fmt(rev26), 10, False, BLACK)
    set_cell(t2, ri, 2, fmt(rev25) if rev25 else '—', 10, False, GRAY if not rev25 else BLACK)
    set_cell(t2, ri, 3, pfmt(rc), 10, False, POS_GREEN if rc>=0 else NEG_RED)
    set_cell(t2, ri, 4, rc_pct, 10, False, BLACK)
    set_cell(t2, ri, 5, fmt(prof26), 10, False, POS_GREEN if prof26>=0 else NEG_RED)
    set_cell(t2, ri, 6, fmt(prof25) if prof25 else '—', 10, False, POS_GREEN if prof25>=0 else NEG_RED if not prof25 else BLACK)
    set_cell(t2, ri, 7, pfmt(pc), 10, False, POS_GREEN if pc>=0 else NEG_RED)
    set_cell(t2, ri, 8, pc_pct, 10, False, BLACK)

rt, pt = R26 - R25, P26 - P25
ti = len(bus_data) + 1
tot_vals = [('一级市场各BU合计', BLACK, PP_ALIGN.LEFT), (fmt(R26), BLACK, PP_ALIGN.CENTER), (fmt(R25), BLACK, PP_ALIGN.CENTER),
            (pfmt(rt), POS_GREEN, PP_ALIGN.CENTER), (f"+{rt/R25*100:.0f}%" if R25 else '—', BLACK, PP_ALIGN.CENTER),
            (fmt(P26), POS_GREEN, PP_ALIGN.CENTER), (fmt(P25), POS_GREEN, PP_ALIGN.CENTER),
            (pfmt(pt), POS_GREEN, PP_ALIGN.CENTER), (f"+{pt/abs(P25)*100:.0f}%" if P25 else '—', BLACK, PP_ALIGN.CENTER)]
for ci, (v, clr, al) in enumerate(tot_vals):
    set_cell(t2, ti, ci, v, 10, True, clr, LIGHT_BG, al)

# RIGHT: Pie charts as images
# Revenue pie
rev_pie_data = [(b[0], b[1]) for b in bus_data]
rev_buf = generate_pie_chart(rev_pie_data, '2026预算收入构成', 'rev_pie')
s2.shapes.add_picture(rev_buf, Inches(8.8), Inches(1.15), Inches(4.1), Inches(2.5))

# Revenue legend
rev_sorted = sorted([(b[0], b[1]) for b in bus_data], key=lambda x: x[1], reverse=True)
rev_legend_text = '  '.join([f'● {n}' for n, v in rev_sorted])
rev_leg = s2.shapes.add_textbox(Inches(8.8), Inches(3.68), Inches(4.1), Inches(0.45))
rev_leg.text_frame.word_wrap = True
rp = rev_leg.text_frame.paragraphs[0]; rp.text = rev_legend_text
rp.font.size = Pt(6); rp.font.color.rgb = GRAY; rp.font.name = 'Microsoft YaHei'

# Profit pie (positive values only)
prof_pie_data = [(b[0], b[3]) for b in bus_data if b[3] > 0]
prof_buf = generate_pie_chart(prof_pie_data, '2026预算利润贡献构成（正值）', 'prof_pie')
s2.shapes.add_picture(prof_buf, Inches(8.8), Inches(4.2), Inches(4.1), Inches(2.5))

# Profit legend
prof_sorted = sorted([(b[0], b[3]) for b in bus_data if b[3] > 0], key=lambda x: x[1], reverse=True)
prof_legend_text = '  '.join([f'● {n}' for n, v in prof_sorted])
prof_leg = s2.shapes.add_textbox(Inches(8.8), Inches(6.72), Inches(4.1), Inches(0.3))
prof_leg.text_frame.word_wrap = True
pp = prof_leg.text_frame.paragraphs[0]; pp.text = prof_legend_text
pp.font.size = Pt(6); pp.font.color.rgb = GRAY; pp.font.name = 'Microsoft YaHei'

fn2 = s2.shapes.add_textbox(Inches(0.5), Inches(6.88), Inches(8), Inches(0.25))
fn2.text_frame.paragraphs[0].text = "注：2025年无数据的BU系当年度尚未设立或未开展运营。"
fn2.text_frame.paragraphs[0].font.size = Pt(9); fn2.text_frame.paragraphs[0].font.color.rgb = GRAY

# ============ SLIDE 3: 费用 ============
s3 = prs.slides.add_slide(prs.slide_layouts[6])
add_header(s3, "2026年度集团预算分析 — 总部职能部门费用", "3")

t3 = s3.shapes.add_table(14, 5, Inches(0.6), Inches(1.15), Inches(12.1), Inches(4.5)).table

h3 = ['费用性质', '2026预算', '2025实际', '变动额', '变动率']
for c, h in enumerate(h3): set_cell(t3, 0, c, h, 11, True, WHITE, DARK_BLUE)

E26 = E25 = 0
for i, (name, v26, v25, note) in enumerate(expense_data):
    E26 += v26; E25 += v25; chg = v26 - v25
    pct = f"{'+' if chg>=0 else ''}{chg/v25*100:.1f}%" if v25 else '—'
    ri = i + 1
    set_cell(t3, ri, 0, name, 11, True, BLACK, None, PP_ALIGN.LEFT)
    set_cell(t3, ri, 1, fmt(v26), 11, False, BLACK)
    set_cell(t3, ri, 2, fmt(v25), 11, False, BLACK)
    set_cell(t3, ri, 3, pfmt(chg), 11, False, NEG_RED if chg>=0 else POS_GREEN)
    set_cell(t3, ri, 4, '注' if note else pct, 11, False, BLACK)

ec = E26 - E25
ti3 = len(expense_data) + 1
for c, v in enumerate(['变动营运费用合计', E26, E25, ec, f"+{ec/E25*100:.1f}%"]):
    set_cell(t3, ti3, c, fmt(v) if isinstance(v, int) else v, 11, True, BLACK, LIGHT_BG, PP_ALIGN.LEFT if c==0 else PP_ALIGN.CENTER)

fn3 = s3.shapes.add_textbox(Inches(0.6), Inches(5.85), Inches(12), Inches(0.9))
fn3.text_frame.word_wrap = True
p3 = fn3.text_frame.paragraphs[0]; p3.font.size = Pt(11); p3.font.color.rgb = RGBColor(0x5A, 0x4E, 0x2F)
p3.text = "注：2025年实际发生的IT及信息系统支出合计人民币308万元，因历史分类方式不同，分散于折旧摊销、专业咨询费及其他业务支出等项目。2026年预算统一归集至资讯系统列示，整体资讯系统费用较2025年实际增加约人民币170万元。"
p3b = fn3.text_frame.add_paragraph(); p3b.font.size = Pt(9); p3b.font.color.rgb = GRAY
p3b.text = "费用数据为总部各职能部门（管委会、HCC、区域平台、财务部、法务部、综合服务中心、投资执行部、投资管理部）合计，不包含策略机遇投资部。"

out = r"D:\Users\yuxr\Desktop\2026集团预算分析.pptx"
prs.save(out)
print(f"PPTX saved: {out}")
print(f"Size: {os.path.getsize(out)} bytes")
